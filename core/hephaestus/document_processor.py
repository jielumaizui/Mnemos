#!/usr/bin/env python3
import logging

logger = logging.getLogger(__name__)

"""
Document Processor - 文档处理器

支持格式：
- Excel (.xlsx, .xls) → Markdown表格
- PPT (.pptx, .ppt) → Markdown幻灯片列表
- PDF (.pdf) → Markdown文本（带页码）
- Word (.docx) → Markdown
- HTML (.html, .htm) → Markdown

默认导入流程：
1. 检测文件类型
2. 提取内容并转为 Markdown
3. DocumentImportService/FileIngestor 写 canonical raw
4. capture outbox → Amphora → 质量门 → Wiki

旧 ``--save``/StorageBackend 直写入口已退役，避免与 raw projection 争夺所有权。
"""

# Document Processor - 文档处理器
#
# 默认导入流程：
# 1. 检测文件类型
# 2. 提取内容并转为 Markdown
# 3. 写 canonical raw
# 4. capture outbox 异步进入质量门与 Wiki
#
# 注意：旧 AgentDelegate 委托验证模式已退役。
# 本模块保留的是当前本地规则验证入口；Mnemos 直接调用 LLM API 执行蒸馏。

import sys  # noqa: E402
import importlib  # noqa: E402
import json  # noqa: E402
import re  # noqa: E402
import shutil  # noqa: E402
import sqlite3  # noqa: E402
import subprocess  # noqa: E402
from datetime import datetime  # noqa: E402
from pathlib import Path  # noqa: E402
from typing import Any, Dict, List, Optional, Tuple, Union  # noqa: E402
from dataclasses import dataclass  # noqa: E402
from enum import Enum  # noqa: E402

sys.path.insert(0, str(Path(__file__).parent.parent.parent))

# Constants extracted from magic numbers
RESULT_SECONDS = 30


DOCUMENT_OPERATION_ERRORS = (
    ImportError,
    KeyError,
    OSError,
    RuntimeError,
    TypeError,
    ValueError,
    sqlite3.Error,
    subprocess.SubprocessError,
)


def _document_processing_errors():
    """Return the concrete optional-library errors handled at the public boundary."""
    error_types = list(DOCUMENT_OPERATION_ERRORS)
    optional_errors = (
        ("bs4.exceptions", "ParserRejectedMarkup"),
        ("docx.opc.exceptions", "OpcError"),
        ("ebooklib.epub", "EpubException"),
        ("openpyxl.utils.exceptions", "InvalidFileException"),
        ("pptx.exc", "PythonPptxError"),
        ("pypdf.errors", "PyPdfError"),
    )
    for module_name, error_name in optional_errors:
        try:
            module = importlib.import_module(module_name)
        except ImportError:
            continue
        error_type = getattr(module, error_name, None)
        if isinstance(error_type, type) and issubclass(error_type, BaseException):
            error_types.append(error_type)
    return tuple(error_types)


class DocumentType(Enum):
    """文档类型"""

    EXCEL = "excel"
    PPT = "ppt"
    PDF = "pdf"
    WORD = "word"
    HTML = "html"
    EBOOK = "ebook"
    UNKNOWN = "unknown"


@dataclass
class ExtractedDocument:
    """提取的文档内容"""

    doc_type: DocumentType
    filename: str
    title: str
    content: str  # Markdown 格式
    metadata: Dict  # 文档元数据
    summary: str  # 内容摘要
    # 验证相关字段
    validation_status: str = "pending"  # pending, validated, review, rejected
    needs_review: bool = False
    review_reason: str = ""
    confidence: float = 0.0  # 提取置信度
    processing_method: str = "local"  # local, cloud, fallback


class DocumentProcessor:
    """
    文档处理器 v2.0

    使用本地规则验证提取质量，并支持人工核对机制。
    自动检测文档类型并提取内容为 Markdown。
    """

    SUPPORTED_EXTENSIONS = {
        ".xlsx",
        ".xls",  # Excel
        ".pptx",  # PowerPoint (.ppt 老格式不直接支持，需另存为 .pptx)
        ".pdf",  # PDF
        ".docx",  # Word (.doc 老格式不直接支持，需另存为 .docx)
        ".html",
        ".htm",  # HTML
        ".epub",
        ".mobi",
        ".azw3",  # Ebook
    }

    # 验证阈值（与ImageProcessor保持一致）
    VALIDATION_CONFIDENCE_THRESHOLD = 0.85  # 验证通过阈值
    REVIEW_CONFIDENCE_THRESHOLD = 0.60  # 人工核对阈值

    def __init__(self):
        self._check_dependencies()

    def _check_dependencies(self):
        """检查并报告依赖可用性"""
        self.deps = {
            "pandas": False,
            "openpyxl": False,
            "pptx": False,
            "pypdf": False,
            "docx": False,
            "beautifulsoup4": False,
            "markdownify": False,
            "ebooklib": False,
        }

        try:
            import pandas  # noqa: F401

            self.deps["pandas"] = True
        except ImportError:
            logger.debug("[document_processor] ImportError suppressed", exc_info=True)

        try:
            import openpyxl  # noqa: F401

            self.deps["openpyxl"] = True
        except ImportError:
            logger.debug("[document_processor] ImportError suppressed", exc_info=True)

        try:
            from pptx import Presentation  # noqa: F401

            self.deps["pptx"] = True
        except ImportError:
            logger.debug("[document_processor] ImportError suppressed", exc_info=True)

        try:
            import pypdf  # noqa: F401

            self.deps["pypdf"] = True
        except ImportError:
            logger.debug("[document_processor] ImportError suppressed", exc_info=True)

        try:
            from docx import Document  # noqa: F401

            self.deps["docx"] = True
        except ImportError:
            logger.debug("[document_processor] ImportError suppressed", exc_info=True)

        try:
            from bs4 import BeautifulSoup  # noqa: F401

            self.deps["beautifulsoup4"] = True
        except ImportError:
            logger.debug("[document_processor] ImportError suppressed", exc_info=True)

        try:
            import markdownify  # noqa: F401

            self.deps["markdownify"] = True
        except ImportError:
            logger.debug("[document_processor] ImportError suppressed", exc_info=True)

        try:
            import ebooklib  # noqa: F401

            self.deps["ebooklib"] = True
        except ImportError:
            logger.debug("[document_processor] ImportError suppressed", exc_info=True)

        # 打印依赖状态（首次初始化时检查，缺失时仅 INFO 提示避免刷屏）
        missing = [k for k, v in self.deps.items() if not v]
        if missing:
            logger.info(
                "[DocumentProcessor] 可选依赖未安装: %s — 对应文档格式处理将不可用",
                ", ".join(missing),
            )

    def _score_content_length(self, content: str) -> Tuple[float, Optional[str]]:
        """内容长度评分（1分）。"""
        if len(content) > 100:
            return 1.0, None
        return 0.0, f"内容过短 ({len(content)} 字符)"

    def _score_title(self, title: str, file_path: Path) -> Tuple[float, Optional[str]]:
        """标题有效性评分（1分）。"""
        if title and title != file_path.stem:
            return 1.0, None
        if title:
            return 0.5, None
        return 0.0, "标题为空"

    def _score_structure(self, content: str) -> Tuple[float, Optional[str]]:
        """Markdown 结构评分（1分）。"""
        has_headers = "#" in content
        has_lists = "- " in content or "* " in content
        has_paragraphs = "\n\n" in content
        structure_score = sum([has_headers, has_lists, has_paragraphs])
        score = min(structure_score / 2, 1.0)
        issue = None if structure_score >= 2 else "结构简单（缺少标题/列表/段落）"
        return score, issue

    def _score_type_match(
        self, doc: ExtractedDocument, file_path: Path
    ) -> Tuple[float, Optional[str]]:
        """文件类型匹配评分（1分）。"""
        expected_type = self.detect_type(file_path)
        if expected_type and doc.doc_type == expected_type:
            return 1.0, None
        return (
            0.0,
            f"类型可能不匹配 (检测到 {doc.doc_type.value}, 期望 {expected_type.value if expected_type else 'unknown'})",  # noqa: E501
        )

    def _score_metadata(self, doc: ExtractedDocument) -> Tuple[float, Optional[str]]:
        """元数据完整性评分（1分）。"""
        meta = doc.metadata or {}
        has_pages = meta.get("pages") or meta.get("sheets") or meta.get("slides")
        has_summary = bool(doc.summary)
        if has_pages and has_summary:
            return 1.0, None
        if has_pages or has_summary:
            return 0.5, None
        return 0.0, "缺少元数据或摘要"

    @staticmethod
    def _validation_outcome(score: float, issues: List[str]) -> Dict:
        """根据总分计算验证结果。"""
        confidence = round(score / 5.0, 2)
        if confidence >= 0.85:
            suggested_action = "pass"
            is_valid = True
        elif confidence >= 0.60:
            suggested_action = "review"
            is_valid = True
        else:
            suggested_action = "review"
            is_valid = False
        return {
            "is_valid": is_valid,
            "confidence": confidence,
            "issues": issues or ["本地规则验证通过"],
            "suggested_action": suggested_action,
        }

    def validate_extraction(self, doc: ExtractedDocument, file_path: Path) -> Dict:
        """
        文档提取结果本地规则验证。

        旧 AgentDelegate 委托验证模式已退役；本方法仍是当前验证入口。

        基于内容质量指标计算置信度：
        - 非空内容、合理长度、结构完整性、标题有效性
        """
        issues: List[str] = []
        score = 0.0

        content = (doc.content or "").strip()
        title = (doc.title or "").strip()

        sub_score, issue = self._score_content_length(content)
        score += sub_score
        if issue:
            issues.append(issue)

        sub_score, issue = self._score_title(title, file_path)
        score += sub_score
        if issue:
            issues.append(issue)

        sub_score, issue = self._score_structure(content)
        score += sub_score
        if issue:
            issues.append(issue)

        sub_score, issue = self._score_type_match(doc, file_path)
        score += sub_score
        if issue:
            issues.append(issue)

        sub_score, issue = self._score_metadata(doc)
        score += sub_score
        if issue:
            issues.append(issue)

        return self._validation_outcome(score, issues)

    def process_document_with_validation(self, file_path: Path) -> Optional[ExtractedDocument]:
        """
        处理文档（带本地规则验证流程）

        完整流程：
        1. 本地提取文档内容
        2. 本地规则验证提取质量
        3. 根据验证结果决定：通过/标记人工审核/拒绝
        """
        logger.info("[DocumentProcessor] 📄 开始处理（带验证）: %s", file_path.name)

        # Step 1: 本地提取
        doc = self.process_document(file_path)
        if not doc:
            return None

        # Step 2: 验证提取结果
        logger.info("[DocumentProcessor] 🔎 使用本地规则验证提取结果...")
        validation = self.validate_extraction(doc, file_path)

        confidence = validation.get("confidence", 0.0)
        is_valid = validation.get("is_valid", False)
        suggested_action = validation.get("suggested_action", "review")
        issues = validation.get("issues", [])

        if suggested_action == "reject":
            # 严重不可信，拒绝入库
            logger.info("[DocumentProcessor] ❌ 验证拒绝 (置信度: %.2f)", confidence)
            logger.info(
                "[DocumentProcessor] 原因: %s", ", ".join(issues) if issues else "内容严重不可信"
            )
            doc.validation_status = "rejected"
            doc.needs_review = False
            doc.confidence = confidence
            doc.review_reason = (
                f"验证拒绝 ({confidence:.2f}): {', '.join(issues) if issues else '内容严重不可信'}"
            )

        elif is_valid and confidence >= self.VALIDATION_CONFIDENCE_THRESHOLD:
            # 验证通过
            logger.info("[DocumentProcessor] ✅ 验证通过 (置信度: %.2f)", confidence)
            doc.validation_status = "validated"
            doc.confidence = confidence
            doc.needs_review = False

        elif suggested_action == "reprocess" or confidence < self.REVIEW_CONFIDENCE_THRESHOLD:
            # 需要重新处理（使用云端）
            logger.warning("[DocumentProcessor] ⚠️ 验证失败/置信度低，需要人工核对...")
            logger.info("[DocumentProcessor] 原因: %s", ", ".join(issues) if issues else "未知")

            # 对于文档，我们无法像图片那样云端重处理
            # 标记为需要人工审核
            doc.validation_status = "review"
            doc.needs_review = True
            doc.confidence = confidence
            doc.review_reason = f"验证置信度低 ({confidence:.2f}): {', '.join(issues) if issues else '结构可能不准确'}"  # noqa: E501

        else:
            # 标记人工核对
            logger.warning("[DocumentProcessor] ⚠️ 标记待人工核对 (置信度: %.2f)", confidence)
            doc.validation_status = "review"
            doc.needs_review = True
            doc.confidence = confidence
            doc.review_reason = f"置信度较低: {', '.join(issues) if issues else '建议人工核对'}"

        return doc

    def detect_type(self, file_path: Path) -> DocumentType:
        """检测文档类型"""
        ext = file_path.suffix.lower()

        if ext in [".xlsx", ".xls"]:
            return DocumentType.EXCEL
        elif ext == ".pptx":
            return DocumentType.PPT
        elif ext == ".pdf":
            return DocumentType.PDF
        elif ext == ".docx":
            return DocumentType.WORD
        elif ext in [".html", ".htm"]:
            return DocumentType.HTML
        elif ext in [".epub", ".mobi", ".azw3"]:
            return DocumentType.EBOOK
        else:
            return DocumentType.UNKNOWN

    def process_document(self, file_path: Path) -> Optional[ExtractedDocument]:
        """
        处理文档

        自动检测类型并提取内容
        """
        if not file_path.exists():
            logger.info("[DocumentProcessor] ❌ 文件不存在: %s", file_path)
            return None

        doc_type = self.detect_type(file_path)

        if doc_type == DocumentType.UNKNOWN:
            logger.warning("[DocumentProcessor] ❌ 不支持的文件类型: %s", file_path.suffix)
            return None

        logger.info("[DocumentProcessor] 📄 处理 %s: %s", doc_type.value, file_path.name)

        # 根据类型处理
        processors = {
            DocumentType.EXCEL: self._process_excel,
            DocumentType.PPT: self._process_ppt,
            DocumentType.PDF: self._process_pdf,
            DocumentType.WORD: self._process_word,
            DocumentType.HTML: self._process_html,
            DocumentType.EBOOK: self._process_ebook,
        }

        processor = processors.get(doc_type)
        if not processor:
            logger.warning("[DocumentProcessor] ⚠️ 暂无 %s 类型处理器", doc_type.value)
            return None

        try:
            doc = processor(file_path)
            if doc is not None:
                doc.metadata["source_path"] = str(file_path)
            return doc
        except _document_processing_errors() as e:
            logger.warning("[DocumentProcessor] ❌ 处理失败: %s", e, exc_info=True)
            return None

    def process_and_distill(
        self,
        file_path: Path,
        inbox: Path | None = None,
        force_provider: str = "api",
        source: str = "human",
        doc: Optional[ExtractedDocument] = None,
        return_details: bool = False,
    ) -> Union[int, Dict[str, Any]]:
        """
        CLI 快速路径：直接处理本地文档并蒸馏入 wiki。

        为保证 provenance 一致，提取后的文档先写 canonical raw revision；该显式
        直出入口声明自己是 distillation owner，使 capture outbox 不再重复投递。

        Args:
            file_path: 本地文件路径
            inbox: wiki Inbox 目录（默认自动检测）
            force_provider: 强制 LLM 提供商 — "api"/None 使用默认 API chain，
                具体 provider 名（如 dmxapi/siliconflow/openai）会收窄 API chain；
                不存在的 provider 会返回明确错误。
            source: Wiki frontmatter / 训练样本中的来源标记
            doc: 已解析出的文档对象；MCP parse 阶段可传入以避免重复解析
            return_details: True 时返回包含 storage_uid/session_id/wiki_paths 的详情字典
        """
        from core.hephaestus.document_pipeline import DocumentDistillationPipeline
        from core.hephaestus.distillation_engine import HttpApiHostAgentCaller
        from core.document_import import file_sha256
        import hashlib

        doc = doc or self.process_document(file_path)
        if not doc:
            if return_details:
                return {
                    "page_count": 0,
                    "wiki_paths": [],
                    "storage_uid": None,
                    "session_id": "",
                    "fragment_count": 0,
                    "judgment": "parse_failed",
                }
            return 0

        from core.privacy.ingestion_security import assess_ingestion_security

        security = assess_ingestion_security(doc.content or "")

        file_hash = hashlib.md5(str(file_path).encode(), usedforsecurity=False).hexdigest()[:8]
        asset_id = f"document:{file_sha256(file_path)[:24]}"
        provisional_session_id = f"doc-{file_hash}"
        from core.sync_framework.ingestion_receipt import create_ingestion_receipt

        ingestion_receipt = create_ingestion_receipt(
            content=doc.content or "",
            source_agent="document_processor",
            source_path=str(file_path),
            session_id=provisional_session_id,
            title=doc.title or doc.filename,
            metadata={
                "source": source,
                "filename": doc.filename,
                "doc_type": doc.doc_type.value,
                "asset_kind": "trusted_user_document",
                "asset_id": asset_id,
                "asset_title": doc.title or doc.filename,
                "distill_requested": False,
                "distillation_owner": "document_pipeline_direct",
                "content_source": "external_file",
                "source_authority": "external_content",
                "source_authority_purpose": "searchable_reference_or_pending_hypothesis",
                "security": security.as_dict(),
            },
        )
        blocked_details = {
            "page_count": 0,
            "wiki_paths": [],
            "storage_uid": None,
            "session_id": provisional_session_id,
            "fragment_count": 0,
            "judgment": "capture_failed_recoverable",
            "ingestion_receipt": ingestion_receipt,
        }
        if not ingestion_receipt.get("success"):
            logger.warning("[DocumentProcessor] canonical capture receipt failed: %s", file_path.name)
            return blocked_details if return_details else 0

        storage_uid = str(ingestion_receipt.get("raw_event_id") or "")
        if not storage_uid:
            logger.warning(
                "[DocumentProcessor] canonical raw revision 缺失，阻断正式蒸馏: %s",
                file_path.name,
            )
            return blocked_details if return_details else 0

        session_id = storage_uid

        # 构建 messages（模拟 reconstruct_session 输出）
        messages = [
            {
                "role": "user",
                "content": doc.content,
                "content_source": "external_file",
                "asset_kind": "trusted_user_document",
                "source_authority": "external_content",
                "source_authority_purpose": "searchable_reference_or_pending_hypothesis",
            }
        ]
        meta = {
            "source": source,
            "filename": doc.filename,
            "file_path": str(file_path),
            "doc_type": doc.doc_type.value,
            "storage_uid": storage_uid,
            "file_hash": file_hash,
            "source_event_id": ingestion_receipt.get("source_event_id", ""),
            "raw_event_id": ingestion_receipt.get("raw_event_id", ""),
            "provenance_id": ingestion_receipt.get("provenance_id", ""),
            "content_source": "external_file",
            "source_authority": "external_content",
            "source_authority_purpose": "searchable_reference_or_pending_hypothesis",
            "security": security.as_dict(),
            "pages": doc.metadata.get(
                "pages", doc.metadata.get("slides", doc.metadata.get("chapters", 0))
            ),
        }

        # 调用文档蒸馏管道（支持强制 provider）
        caller = HttpApiHostAgentCaller(force_provider=force_provider)
        pipeline = DocumentDistillationPipeline(caller=caller)
        if inbox:
            pipeline.inbox_dir = inbox

        result = pipeline.process(session_id, messages, meta)

        if result.judgment == "skip" or not result.fragments:
            logger.info("[DocumentProcessor] 文档被跳过: %s", file_path.name)
            if return_details:
                return {
                    "page_count": 0,
                    "wiki_paths": [],
                    "storage_uid": storage_uid,
                    "session_id": session_id,
                    "fragment_count": 0,
                    "judgment": result.judgment,
                    "doc_category": getattr(result, "doc_category", ""),
                }
            return 0

        # 写入 wiki
        paths = pipeline.write_to_wiki(result, source=source)

        # 写入画像信号
        try:
            from core.persona.psyche import get_signal_store
            from datetime import datetime

            store = get_signal_store()
            store.insert_document_signal(
                session_id=session_id,
                filename=doc.filename,
                doc_type=doc.doc_type.value,
                doc_category=result.doc_category,
                title=doc.title,
                key_topics=json.dumps(
                    result.fragments[0].keywords if result.fragments else [], ensure_ascii=False
                ),
                entity_type=(
                    result.fragments[0].frontmatter.get("类型", "reference")
                    if result.fragments
                    else "reference"
                ),
                page_count=meta.get("pages", 0),
                import_timestamp=datetime.now().isoformat(),
                import_source=str(file_path),
                confidence=0.8,
            )
        except DOCUMENT_OPERATION_ERRORS as e:
            logger.debug("[DocumentProcessor] 文档信号写入失败: %s", e)

        logger.info(
            "[DocumentProcessor] ✅ 文档已蒸馏入 wiki: %s → %s 页", file_path.name, len(paths)
        )
        if return_details:
            return {
                "page_count": len(paths),
                "wiki_paths": paths,
                "storage_uid": storage_uid,
                "session_id": session_id,
                "fragment_count": len(result.fragments) if result and result.fragments else 0,
                "judgment": result.judgment,
                "doc_category": getattr(result, "doc_category", ""),
                "ingestion_receipt": ingestion_receipt,
            }
        return len(paths)

    def _process_excel(self, file_path: Path) -> Optional[ExtractedDocument]:
        """
        处理 Excel 文件

        提取所有工作表为 Markdown 表格
        """
        if not self.deps["pandas"] or not self.deps["openpyxl"]:
            # 回退：使用系统命令转换为 CSV 再处理
            return self._process_excel_fallback(file_path)

        import pandas as pd

        # 读取所有工作表
        xl_file = pd.ExcelFile(file_path)
        sheet_names = xl_file.sheet_names

        content_lines = [f"# 📊 Excel: {file_path.stem}", ""]
        metadata = {
            "sheets": len(sheet_names),
            "sheet_names": sheet_names,
            "total_rows": 0,
            "total_cells": 0,
        }

        for sheet_name in sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)

            # 跳过空表
            if df.empty:
                continue

            content_lines.append(f"## 工作表: {sheet_name}")
            content_lines.append("")

            # 转换为 Markdown 表格
            # 保留全部行；canonical raw 与后续 token 分块负责完整内容。
            display_df = df

            # 生成表头
            headers = [str(col) for col in display_df.columns]
            content_lines.append("| " + " | ".join(headers) + " |")
            content_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

            # 生成数据行
            for _, row in display_df.iterrows():
                cells = [str(cell) if pd.notna(cell) else "" for cell in row]
                # 截断过长的单元格（防止表格变形）
                cells = [c[:200] + "..." if len(c) > 200 else c for c in cells]
                content_lines.append("| " + " | ".join(cells) + " |")

            content_lines.append("")

            metadata["total_rows"] += len(df)
            metadata["total_cells"] += len(df) * len(df.columns)

        summary = f"Excel文件：{len(sheet_names)}个工作表，共{metadata['total_rows']}行"

        return ExtractedDocument(
            doc_type=DocumentType.EXCEL,
            filename=file_path.name,
            title=file_path.stem,
            content="\n".join(content_lines),
            metadata=metadata,
            summary=summary,
            validation_status="pending",
            needs_review=False,
            review_reason="",
            confidence=0.8,  # Excel提取通常较准确
            processing_method="local",
        )

    def _process_excel_fallback(self, file_path: Path) -> Optional[ExtractedDocument]:
        """Excel 处理回退方案（使用系统命令）"""
        logger.warning("[DocumentProcessor] ⚠️ 使用回退方案处理 Excel...")

        import tempfile

        tmp_dir = Path(tempfile.gettempdir())

        # 尝试使用 ssconvert (Gnumeric) 或 LibreOffice 转换
        tmp_dir / f"excel_export_{datetime.now().strftime('%Y%m%d%H%M%S')}.csv"

        try:
            # 检查 LibreOffice 是否可用
            if not shutil.which("libreoffice"):
                raise FileNotFoundError("libreoffice 未安装")

            # 尝试使用 LibreOffice 转换
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "csv",
                    "--outdir",
                    str(tmp_dir),
                    str(file_path),
                ],
                capture_output=True,
                timeout=RESULT_SECONDS,
            )

            if result.returncode == 0:
                # 查找生成的 CSV 文件
                csv_files = list(tmp_dir.glob(f"{file_path.stem}*.csv"))
                if csv_files:
                    # 读取 CSV 并转为 Markdown
                    content_lines = [f"# 📊 Excel: {file_path.stem}", ""]

                    with open(csv_files[0], "r", encoding="utf-8", errors="ignore") as f:
                        lines = f.readlines()

                    if lines:
                        # 第一行作为表头
                        headers = lines[0].strip().split(",")
                        content_lines.append("| " + " | ".join(headers) + " |")
                        content_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

                        # 数据行（限制前100行）
                        for line in lines[1:101]:
                            cells = line.strip().split(",")
                            content_lines.append("| " + " | ".join(cells) + " |")

                        if len(lines) > 101:
                            content_lines.append(f"\n*... 共 {len(lines)-1} 行，显示前 100 行*")

                    # 清理临时文件
                    for cf in csv_files:
                        cf.unlink()

                    return ExtractedDocument(
                        doc_type=DocumentType.EXCEL,
                        filename=file_path.name,
                        title=file_path.stem,
                        content="\n".join(content_lines),
                        metadata={"method": "libreoffice_fallback"},
                        summary="Excel文件（LibreOffice转换）",
                        validation_status="pending",
                        needs_review=False,
                        review_reason="",
                        confidence=0.6,  # 回退方案置信度较低
                        processing_method="fallback",
                    )

        except (OSError, TypeError, ValueError, subprocess.SubprocessError) as e:
            logger.warning("[DocumentProcessor] ❌ 回退处理失败: %s", e)

        return None

    def _process_ppt(self, file_path: Path) -> Optional[ExtractedDocument]:
        """
        处理 PowerPoint 文件

        提取每页的标题和内容
        """
        if not self.deps["pptx"]:
            logger.info("[DocumentProcessor] ❌ 缺少 python-pptx，无法处理 PPT")
            logger.info("[DocumentProcessor] 安装: pip install python-pptx")
            return None

        from pptx import Presentation

        prs = Presentation(file_path)  # type: ignore[arg-type]

        content_lines = [f"# 📽️ PowerPoint: {file_path.stem}", ""]
        content_lines.append(f"**幻灯片数量**: {len(prs.slides)}")
        content_lines.append("")

        slides_content = []
        total_text_chars = 0

        for i, slide in enumerate(prs.slides, 1):
            slide_lines = [f"## 幻灯片 {i}", ""]

            # 提取所有文本
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                    total_text_chars += len(shape.text)

            if slide_texts:
                # 第一行通常作为标题
                slide_lines.append(f"### {slide_texts[0]}")
                slide_lines.append("")

                # 其余作为内容
                for text in slide_texts[1:]:
                    # 分段处理
                    paragraphs = text.split("\n")
                    for para in paragraphs:
                        if para.strip():
                            slide_lines.append(para.strip())

            slides_content.append("\n".join(slide_lines))

        # 合并所有幻灯片内容
        content_lines.extend(slides_content)

        metadata = {"slides": len(prs.slides), "total_text_chars": total_text_chars}

        summary = f"PPT文件：{len(prs.slides)}页幻灯片"

        return ExtractedDocument(
            doc_type=DocumentType.PPT,
            filename=file_path.name,
            title=file_path.stem,
            content="\n\n".join(content_lines),
            metadata=metadata,
            summary=summary,
            validation_status="pending",
            needs_review=False,
            review_reason="",
            confidence=0.75,
            processing_method="local",
        )

    def _process_pdf(self, file_path: Path) -> Optional[ExtractedDocument]:
        """
        处理 PDF 文件

        提取文本并保留页码信息
        """
        if not self.deps["pypdf"]:
            # 回退到 pdftotext (poppler)
            return self._process_pdf_fallback(file_path)

        import pypdf

        pdf_error_type = getattr(getattr(pypdf, "errors", None), "PyPdfError", OSError)

        content_lines = [f"# 📄 PDF: {file_path.stem}", ""]
        metadata = {"pages": 0, "extracted_pages": 0}

        try:
            with open(file_path, "rb") as f:
                pdf_reader = pypdf.PdfReader(f)
                num_pages = len(pdf_reader.pages)
                metadata["pages"] = num_pages

                content_lines.append(f"**总页数**: {num_pages}")
                content_lines.append("")

                # 提取每页内容
                for i, page in enumerate(pdf_reader.pages, 1):
                    try:
                        text = page.extract_text()
                        if text and text.strip():
                            content_lines.append(f"## 第 {i} 页")
                            content_lines.append("")
                            content_lines.append(text.strip())
                            content_lines.append("")
                            metadata["extracted_pages"] += 1
                    except (pdf_error_type, TypeError, ValueError) as e:
                        content_lines.append(f"*第 {i} 页提取失败: {e}*")
                        content_lines.append("")

        except (OSError, pdf_error_type, TypeError, ValueError) as e:
            logger.warning("[DocumentProcessor] ❌ PDF 处理失败: %s", e, exc_info=True)
            return self._process_pdf_fallback(file_path)

        summary = f"PDF文件：{metadata['pages']}页，成功提取{metadata['extracted_pages']}页"

        return ExtractedDocument(
            doc_type=DocumentType.PDF,
            filename=file_path.name,
            title=file_path.stem,
            content="\n".join(content_lines),
            metadata=metadata,
            summary=summary,
            validation_status="pending",
            needs_review=False,
            review_reason="",
            confidence=0.7,
            processing_method="local",
        )

    def _process_pdf_fallback(self, file_path: Path) -> Optional[ExtractedDocument]:
        """PDF 处理回退方案（使用 pdftotext）"""
        logger.warning("[DocumentProcessor] ⚠️ 使用回退方案处理 PDF...")

        # 检查 pdftotext 是否可用
        if not shutil.which("pdftotext"):
            logger.warning("[DocumentProcessor] ⚠️ pdftotext 未安装，跳过 PDF 回退处理")
            return None

        try:
            result = subprocess.run(
                ["pdftotext", "-layout", str(file_path), "-"],
                capture_output=True,
                text=True,
                timeout=60,
            )

            if result.returncode == 0 and result.stdout:
                text = result.stdout

                # 尝试按页分割（如果 pdftotext 支持 -f -l 参数）
                content_lines = [f"# 📄 PDF: {file_path.stem}", ""]
                content_lines.append(text)  # 完整保留

                return ExtractedDocument(
                    doc_type=DocumentType.PDF,
                    filename=file_path.name,
                    title=file_path.stem,
                    content="\n".join(content_lines),
                    metadata={"method": "pdftotext_fallback"},
                    summary="PDF文件（pdftotext提取）",
                    validation_status="pending",
                    needs_review=False,
                    review_reason="",
                    confidence=0.5,  # 回退方案置信度低
                    processing_method="fallback",
                )

        except (OSError, TypeError, ValueError, subprocess.SubprocessError) as e:
            logger.warning("[DocumentProcessor] ❌ 回退处理失败: %s", e)

        return None

    def _process_word(self, file_path: Path) -> Optional[ExtractedDocument]:
        """
        处理 Word 文件

        提取段落和表格
        """
        if not self.deps["docx"]:
            logger.info("[DocumentProcessor] ❌ 缺少 python-docx，无法处理 Word")
            logger.info("[DocumentProcessor] 安装: pip install python-docx")
            return None

        from docx import Document

        doc = Document(file_path)  # type: ignore[arg-type]

        content_lines = [f"# 📝 Word: {file_path.stem}", ""]
        metadata = {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}

        # 提取段落
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # 根据样式判断标题级别
                if para.style.name.startswith("Heading"):  # type: ignore[union-attr]
                    level = para.style.name.replace("Heading ", "")  # type: ignore[union-attr]
                    try:
                        level_num = int(level)
                        content_lines.append(f"{'#' * level_num} {text}")
                    except ValueError:
                        content_lines.append(f"## {text}")
                else:
                    content_lines.append(text)

        # 提取表格
        if doc.tables:
            content_lines.append("")
            content_lines.append("## 表格")
            content_lines.append("")

            for i, table in enumerate(doc.tables, 1):
                content_lines.append(f"### 表格 {i}")
                content_lines.append("")

                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    content_lines.append("| " + " | ".join(cells) + " |")

                content_lines.append("")

        summary = f"Word文件：{metadata['paragraphs']}段落，{metadata['tables']}表格"

        return ExtractedDocument(
            doc_type=DocumentType.WORD,
            filename=file_path.name,
            title=file_path.stem,
            content="\n\n".join(content_lines),
            metadata=metadata,
            summary=summary,
            validation_status="pending",
            needs_review=False,
            review_reason="",
            confidence=0.75,
            processing_method="local",
        )

    def _process_html(self, file_path: Path) -> Optional[ExtractedDocument]:
        """
        处理 HTML 文件

        转换为 Markdown
        """
        content = file_path.read_text(encoding="utf-8", errors="ignore")

        if self.deps["markdownify"]:
            import markdownify

            markdown_content = markdownify.markdownify(content, heading_style="ATX")
        elif self.deps["beautifulsoup4"]:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(content, "html.parser")
            # 移除 script 和 style
            for script in soup(["script", "style"]):
                script.decompose()
            markdown_content = soup.get_text(separator="\n", strip=True)
        else:
            # 纯文本提取
            # 移除 HTML 标签的简单实现
            markdown_content = re.sub(r"<[^>]+>", "", content)
            markdown_content = re.sub(r"\n\s*\n", "\n\n", markdown_content)

        # 添加标题
        title = file_path.stem
        # 尝试从 HTML 中提取 title
        title_match = re.search(r"<title[^>]*>([^<]+)</title>", content, re.IGNORECASE)
        if title_match:
            title = title_match.group(1).strip()

        content_lines = [
            f"# 🌐 HTML: {title}",
            "",
            f"**源文件**: {file_path.name}",
            "",
            markdown_content,
        ]

        metadata = {"original_size": len(content), "extracted_size": len(markdown_content)}

        summary = f"HTML文件：原始{len(content)}字符，提取{len(markdown_content)}字符"

        return ExtractedDocument(
            doc_type=DocumentType.HTML,
            filename=file_path.name,
            title=title,
            content="\n".join(content_lines),
            metadata=metadata,
            summary=summary,
            validation_status="pending",
            needs_review=False,
            review_reason="",
            confidence=0.7,
            processing_method="local",
        )

    def _process_ebook(self, file_path: Path) -> Optional[ExtractedDocument]:
        """
        处理电子书文件（epub/mobi/azw3）

        提取章节结构和文本内容
        """
        ext = file_path.suffix.lower()
        if ext == ".epub":
            return self._process_epub(file_path)
        # mobi/azw3 暂不支持，返回基础信息
        logger.warning("[DocumentProcessor] ⚠️ %s 格式暂不支持完整提取，尝试作为纯文本读取", ext)
        try:
            content = file_path.read_text(encoding="utf-8", errors="ignore")
            return ExtractedDocument(
                doc_type=DocumentType.EBOOK,
                filename=file_path.name,
                title=file_path.stem,
                content=f"# 📖 Ebook: {file_path.stem}\n\n**格式**: {ext}\n\n{content[:5000]}",
                metadata={"format": ext, "size": len(content)},
                summary=f"电子书文件：{ext}格式",
                validation_status="review",
                needs_review=True,
                review_reason=f"{ext}格式暂不支持完整提取，仅保留前5000字符",
                confidence=0.3,
                processing_method="fallback",
            )
        except (OSError, TypeError, ValueError) as e:
            logger.warning("[DocumentProcessor] ❌ 电子书处理失败: %s", e)
            return None

    def _process_epub(self, file_path: Path) -> Optional[ExtractedDocument]:
        """处理 EPUB 电子书"""
        try:
            import ebooklib
            from ebooklib import epub
            from ebooklib.epub import EpubException
        except ImportError:
            logger.warning("[DocumentProcessor] ⚠️ 缺少 ebooklib，无法处理 EPUB")
            logger.info("[DocumentProcessor] 安装: pip install EbookLib")
            return None

        try:
            book = epub.read_epub(str(file_path))

            # 提取元数据
            title = file_path.stem
            try:
                titles = book.get_metadata("DC", "title")
                if titles:
                    title = titles[0][0]
            except (
                OSError, ValueError, TypeError, KeyError, ImportError, AttributeError, RuntimeError,
                subprocess.SubprocessError
            ):
                logger.debug(
                    "[DocumentProcessor] EPUB 元数据提取失败，使用文件名: %s", file_path.name
                )

            # 提取章节内容
            content_lines = [f"# 📖 Ebook: {title}", ""]
            chapters = []
            chapter_count = 0

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    chapter_count += 1
                    try:
                        from bs4 import BeautifulSoup
                        from bs4.exceptions import ParserRejectedMarkup

                        soup = BeautifulSoup(item.get_content(), "html.parser")
                        # 移除 script/style
                        for tag in soup(["script", "style"]):
                            tag.decompose()
                        text = soup.get_text(separator="\n", strip=True)
                        if text.strip():
                            chapters.append(f"## 第 {chapter_count} 章\n\n{text.strip()}")
                    except (OSError, ParserRejectedMarkup, TypeError, ValueError) as e:
                        logger.debug("[DocumentProcessor] BeautifulSoup 解析失败，尝试回退: %s", e)
                        # 回退：直接提取文本
                        try:
                            raw_content = item.get_content()
                            if not isinstance(raw_content, (bytes, bytearray)):
                                raise TypeError("EPUB chapter content must be bytes")
                            text = raw_content.decode("utf-8", errors="ignore")
                            text = re.sub(r"<[^>]+>", "", text)
                            text = re.sub(r"\n\s*\n", "\n\n", text)
                            if text.strip():
                                chapters.append(f"## 第 {chapter_count} 章\n\n{text.strip()}")
                        except (OSError, TypeError, UnicodeError, ValueError) as e2:
                            logger.debug("[DocumentProcessor] EPUB 章节文本提取失败，跳过: %s", e2)
                            continue

            content_lines.extend(chapters)

            metadata = {"format": "epub", "chapters": chapter_count, "title": title}

            summary = f"EPUB电子书：{title}，共{chapter_count}章"

            return ExtractedDocument(
                doc_type=DocumentType.EBOOK,
                filename=file_path.name,
                title=title,
                content="\n\n".join(content_lines),
                metadata=metadata,
                summary=summary,
                validation_status="pending",
                needs_review=False,
                review_reason="",
                confidence=0.7,
                processing_method="local",
            )

        except (OSError, EpubException, TypeError, ValueError, KeyError) as e:
            logger.warning("[DocumentProcessor] ❌ EPUB 处理失败: %s", e, exc_info=True)
            return None

    def save_to_rejected(self, doc: ExtractedDocument, file_path: Path) -> Path:
        """
        将验证拒绝的文档保存到隔离目录，不入主库
        返回保存路径
        """
        from core.config import get_config

        rejected_dir = get_config().database_dir / "rejected_documents"
        rejected_dir.mkdir(parents=True, exist_ok=True)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        base_name = f"{timestamp}_{doc.filename}"

        # 保存元数据
        meta = {
            "filename": doc.filename,
            "doc_type": doc.doc_type.value,
            "validation_status": doc.validation_status,
            "confidence": doc.confidence,
            "review_reason": doc.review_reason,
            "original_path": str(file_path),
            "rejected_at": datetime.now().isoformat(),
            "content_preview": doc.content[:2000],
        }
        meta_path = rejected_dir / f"{base_name}.json"
        meta_path.write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")

        # 保存原始文件副本
        if file_path.exists():
            import shutil

            dest = rejected_dir / f"{base_name}_orig{file_path.suffix}"
            shutil.copy2(file_path, dest)

        logger.info("[DocumentProcessor] 🚫 已拒绝并隔离: %s", meta_path)
        return meta_path


def main():
    """CLI入口"""
    import argparse

    parser = argparse.ArgumentParser(description="Document Processor")
    parser.add_argument("file", nargs="?", help="文档文件路径")
    parser.add_argument("--check-deps", action="store_true", help="检查依赖")

    args = parser.parse_args()

    processor = DocumentProcessor()

    if args.check_deps:
        logger.info("📦 依赖状态:")
        for dep, available in processor.deps.items():
            status = "✅" if available else "❌"
            logger.info("  %s %s", status, dep)
        return

    if not args.file:
        parser.print_help()
        return

    file_path = Path(args.file)
    doc = processor.process_document(file_path)

    if doc:
        logger.info("\n%s", "=" * 50)
        logger.info("处理结果:")
        logger.info("  类型: %s", doc.doc_type.value)
        logger.info("  标题: %s", doc.title)
        logger.info("  摘要: %s", doc.summary)
        logger.info("%s", "=" * 50)
        logger.info("\n内容预览（前500字符）:")
        logger.info(doc.content[:500])
        logger.info("\n... (%s 字符)", len(doc.content))

    else:
        logger.warning("❌ 处理失败")


if __name__ == "__main__":
    main()
