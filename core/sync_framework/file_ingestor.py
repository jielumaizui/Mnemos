# -*- coding: utf-8 -*-
"""
FileIngestor — 用户文件摄入器

将用户文件（PDF/Word/PPT/Excel/HTML/epub/txt/md）完整提取为文本，
写入 canonical raw event store，并把后续投影/蒸馏交给 capture outbox。

设计原则：
  - 文本提取：纯工具提取，零 LLM 成本
  - 大文件处理：复用 document_process.max_file_size_mb，canonical raw 完整保存
  - 编码回退：utf-8 → gbk → latin-1
  - 单一写入所有者：不直接写 StorageBackend，不直接投递 Amphora
"""

from __future__ import annotations

import hashlib
import logging
import os
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

from core.config import ConfigProvider, get_config
from core.document_import import (
    BLOCKED_TEMP_PREFIXES,
    file_sha256,
    validate_trusted_user_document,
)
from core.sync_framework.storage_backend import StorageResult

# Constants extracted from magic numbers
RESULT_SECONDS = 30

logger = logging.getLogger(__name__)

# 明确禁止的系统级临时目录前缀（小写，用于大小写不敏感比较）
_BLOCKED_TEMP_PREFIXES = BLOCKED_TEMP_PREFIXES

# 文件扩展名到提取器的映射
_TEXT_EXTENSIONS = {".md", ".txt", ".csv", ".json", ".yaml", ".yml", ".toml", ".ini", ".log"}
_PDF_EXTENSIONS = {".pdf"}
_DOCX_EXTENSIONS = {".docx"}  # .doc 老格式不直接支持，需另存为 .docx
_PPTX_EXTENSIONS = {".pptx"}  # .ppt 老格式不直接支持，需另存为 .pptx
_XLSX_EXTENSIONS = {".xlsx", ".xls"}
_HTML_EXTENSIONS = {".html", ".htm"}
_EPUB_EXTENSIONS = {".epub"}


class FileIngestor:
    """用户文件摄入器"""

    def __init__(
        self,
        config: Optional[ConfigProvider] = None,
        receipt_factory: Optional[Callable[..., Dict[str, Any]]] = None,
    ):
        self.config = config or get_config()
        self._receipt_factory = receipt_factory
        self.last_session_id: Optional[str] = None
        self.last_queue_id: Optional[str] = None
        self.last_handoff_status: Optional[str] = None
        self.last_projection_status: Optional[str] = None
        self.last_security_assessment: Optional[Dict[str, Any]] = None
        self.last_ingestion_receipt: Optional[Dict[str, Any]] = None

    def _validate_file_path(self, file_path: Path) -> Optional[str]:
        """
        校验摄入路径安全性。

        返回：None 表示通过；否则返回拒绝原因字符串。
        """
        validation = validate_trusted_user_document(
            file_path,
            config=self.config,
            blocked_temp_prefixes=_BLOCKED_TEMP_PREFIXES,
        )
        return None if validation.ok else validation.reason

    def ingest_file(
        self,
        file_path: Path,
        agent_name: str = "file",
        request_distillation: bool = True,
        title: str = "",
    ) -> Optional[List[StorageResult]]:
        """
        摄入单个文件：提取文本 → 构建 Markdown → 写 canonical raw

        Args:
            file_path: 文件路径
            agent_name: 来源 Agent 名
            request_distillation: 是否由 capture outbox 异步请求蒸馏
            title: 文档资产标题；为空时使用文件名

        Returns:
            canonical raw revision receipt 列表，失败返回 None
        """
        file_path = Path(file_path)
        reason = self._validate_file_path(file_path)
        if reason:
            logger.warning("[FileIngestor] %s: %s", reason, file_path)
            return None

        # 文本提取
        text = self._extract_text(file_path)
        if not text:
            logger.warning("[FileIngestor] 无法提取文本: %s", file_path)
            return None

        from core.privacy.ingestion_security import assess_ingestion_security, merge_security_tags

        security = assess_ingestion_security(text)
        self.last_security_assessment = security.as_dict()

        # 不截断内容；canonical raw 必须保存完整的用户可见文本。

        # 构建 Markdown 内容
        content = self._build_file_markdown(file_path, text)
        receipt_title = title or f"file-{file_path.stem}"
        asset_id = f"document:{file_sha256(file_path)[:24]}"

        # 为每个文件生成唯一 session_id，避免多个文件堆进同一个空 session
        file_hash = hashlib.md5(str(file_path).encode("utf-8"), usedforsecurity=False).hexdigest()[:8]
        file_session_id = f"file:{file_hash}:{self._now_date()}"

        # 标签
        tags = [
            f"source={agent_name}",
            f"session={file_session_id}",
            f"time={self._now_date()}",
            "model=file-ingestor",
            # Imported files are external content until an explicit consent
            # flow assigns an authorized project/session ACL.
            "scope=restricted",
            "status=raw",
            "content_type=file-extract",
            "layer=L1",
            f"original-path={file_path.name}",
            f"file-ext={file_path.suffix.lstrip('.')}",
        ]
        tags = merge_security_tags(tags, security)
        receipt = self._create_ingestion_receipt(
            content=content,
            source_agent=f"file_ingestor:{agent_name}",
            source_path=str(file_path),
            session_id=file_session_id,
            title=receipt_title,
            metadata={
                "tags": tags,
                "security": self.last_security_assessment or {},
                "asset_kind": "trusted_user_document",
                "asset_id": asset_id,
                "asset_title": receipt_title,
                "distill_requested": request_distillation,
                "content_source": "external_file",
                "source_authority": "external_content",
                "source_authority_purpose": "searchable_reference_or_pending_hypothesis",
            },
        )
        self.last_ingestion_receipt = receipt
        raw_revision_id = str(receipt.get("raw_event_id") or "")
        if not receipt.get("success") or not raw_revision_id:
            logger.warning("[FileIngestor] canonical capture receipt failed: %s", file_path)
            return None
        tags.extend(
            [
                f"source_event_id={receipt.get('source_event_id', '')}",
                f"raw_event_id={receipt.get('raw_event_id', '')}",
            ]
        )

        self.last_session_id = file_session_id
        capture_result = receipt.get("capture_result") or {}
        self.last_queue_id = str(capture_result.get("capture_dedupe_key") or raw_revision_id)
        self.last_handoff_status = (
            "existing"
            if request_distillation and receipt.get("status") == "duplicate"
            else ("pending" if request_distillation else "not_requested")
        )
        self.last_projection_status = "pending"
        return [
            StorageResult(
                uid=raw_revision_id,
                content=content,
                tags=tags,
                metadata={
                    "canonical_owner": "raw_event_store",
                    "raw_revision_id": raw_revision_id,
                    "capture_queue_ref": self.last_queue_id,
                    "handoff_status": self.last_handoff_status,
                    "projection_status": self.last_projection_status,
                    "asset_kind": "trusted_user_document",
                    "asset_id": asset_id,
                    "asset_title": receipt_title,
                    "content_source": "external_file",
                    "source_authority": "external_content",
                    "source_authority_purpose": (
                        "searchable_reference_or_pending_hypothesis"
                    ),
                },
            )
        ]

    def _create_ingestion_receipt(self, **kwargs) -> Dict[str, Any]:
        if self._receipt_factory is not None:
            return self._receipt_factory(**kwargs)
        from core.sync_framework.ingestion_receipt import create_ingestion_receipt

        return create_ingestion_receipt(**kwargs)

    def ingest_directory(
        self, dir_path: Path, agent_name: str = "file", recursive: bool = True
    ) -> int:
        """
        批量摄入目录中的文件

        Returns:
            成功摄入的文件数量
        """
        if not dir_path.exists() or not dir_path.is_dir():
            return 0
        if dir_path.is_symlink():
            logger.warning("[FileIngestor] 拒绝摄入符号链接目录: %s", dir_path)
            return 0

        count = 0
        pattern = "**/*" if recursive else "*"
        for f in dir_path.glob(pattern):
            if not f.is_file():
                continue
            if not self._is_supported(f):
                continue
            reason = self._validate_file_path(f)
            if reason:
                logger.debug("[FileIngestor] 跳过 %s: %s", f, reason)
                continue
            result = self.ingest_file(f, agent_name)
            if result:
                count += 1
        return count

    def _extract_text(self, file_path: Path) -> Optional[str]:
        """根据文件类型提取文本"""
        ext = file_path.suffix.lower()

        if ext in _TEXT_EXTENSIONS:
            return self._extract_plain(file_path)
        elif ext in _PDF_EXTENSIONS:
            return self._extract_pdf(file_path)
        elif ext in _DOCX_EXTENSIONS:
            return self._extract_docx(file_path)
        elif ext in _PPTX_EXTENSIONS:
            return self._extract_pptx(file_path)
        elif ext in _XLSX_EXTENSIONS:
            return self._extract_xlsx(file_path)
        elif ext in _HTML_EXTENSIONS:
            return self._extract_html(file_path)
        elif ext in _EPUB_EXTENSIONS:
            return self._extract_epub(file_path)
        else:
            logger.debug("[FileIngestor] 不支持的文件类型: %s", ext)
            return None

    def _extract_plain(self, file_path: Path) -> Optional[str]:
        """提取纯文本文件（编码回退），使用 O_NOFOLLOW 防止 TOCTOU 符号链接替换。"""
        flags = os.O_RDONLY | os.O_NOFOLLOW
        if hasattr(os, "O_CLOEXEC"):
            flags |= os.O_CLOEXEC
        try:
            fd = os.open(str(file_path), flags)
        except (OSError, IOError) as e:
            logger.warning("[FileIngestor] 无法打开文件 %s: %s", file_path, e)
            return None
        try:
            for encoding in ("utf-8", "gbk", "latin-1"):
                try:
                    os.lseek(fd, 0, os.SEEK_SET)
                    with open(fd, "r", encoding=encoding, closefd=False) as f:
                        return f.read()
                except (UnicodeDecodeError, UnicodeError):
                    continue
        except (OSError, IOError):
            logger.warning("[FileIngestor] 读取文件失败: %s", file_path, exc_info=True)
            return None
        finally:
            os.close(fd)
        return None

    def _extract_pdf(self, file_path: Path) -> Optional[str]:
        """提取 PDF 文本（完整提取，大内容由 StorageBackend.save() 自动处理）"""
        try:
            import pdfplumber
            from pdfminer.pdfexceptions import PDFException
        except ImportError:
            logger.debug("[FileIngestor] pdfplumber 未安装，尝试 pypdf")
        else:
            extraction_errors: tuple[type[BaseException], ...] = (
                OSError,
                PDFException,
                ValueError,
                TypeError,
            )
            try:
                from pdfplumber.utils.exceptions import PdfminerException
            except ImportError:
                pass
            else:
                extraction_errors = (*extraction_errors, PdfminerException)
            try:
                with pdfplumber.open(str(file_path)) as pdf:
                    texts = []
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            texts.append(text)
                    return "\n\n".join(texts)
            except extraction_errors as e:
                logger.warning(
                    "[FileIngestor] pdfplumber PDF 提取失败，尝试 pypdf: %s",
                    e,
                    exc_info=True,
                )

        pdf_text = self._extract_pdf_pypdf(file_path)
        if pdf_text:
            return pdf_text
        return self._extract_pdf_fallback(file_path)

    def _extract_pdf_pypdf(self, file_path: Path) -> Optional[str]:
        """pypdf 回退：项目依赖已声明 pypdf，优先于系统 pdftotext。"""
        try:
            import pypdf
        except ImportError:
            logger.debug("[FileIngestor] pypdf 未安装，尝试 pdftotext")
            return None
        pypdf_errors = getattr(pypdf, "errors", None)
        pdf_error_type = getattr(pypdf_errors, "PyPdfError", OSError)

        try:
            texts = []
            with open(file_path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        texts.append(text)
            return "\n\n".join(texts) if texts else None
        except (OSError, pdf_error_type, ValueError, TypeError) as e:
            logger.warning(
                "[FileIngestor] pypdf PDF 提取失败，尝试 pdftotext: %s", e, exc_info=True
            )
        return None

    def _extract_pdf_fallback(self, file_path: Path) -> Optional[str]:
        """pdftotext 回退"""
        import subprocess

        try:
            result = subprocess.run(
                ["pdftotext", str(file_path), "-"],
                capture_output=True,
                text=True,
                timeout=RESULT_SECONDS,
            )
            if result.returncode == 0 and result.stdout:
                return result.stdout
        except (FileNotFoundError, subprocess.TimeoutExpired):
            logger.warning(
                "[file_ingestor] (FileNotFoundError, subprocess.TimeoutExpired) suppressed",
                exc_info=True,
            )
        return None

    def _extract_docx(self, file_path: Path) -> Optional[str]:
        """提取 Word 文档文本"""
        try:
            from docx import Document
            from docx.opc.exceptions import OpcError
        except ImportError:
            logger.debug("[FileIngestor] python-docx 未安装")
            return None

        try:
            doc = Document(str(file_path))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text)
        except (OSError, OpcError, ValueError, TypeError, KeyError) as e:
            logger.warning("[FileIngestor] DOCX 提取失败: %s", e, exc_info=True)
            return None

    def _extract_pptx(self, file_path: Path) -> Optional[str]:
        """提取 PPT 文本"""
        try:
            from pptx import Presentation
            from pptx.exc import PythonPptxError
        except ImportError:
            logger.debug("[FileIngestor] python-pptx 未安装")
            return None

        try:
            prs = Presentation(str(file_path))
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text:
                                texts.append(para.text)
                if texts:
                    slides.append(f"## Slide {i+1}\n\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except (OSError, PythonPptxError, ValueError, TypeError, KeyError) as e:
            logger.warning("[FileIngestor] PPTX 提取失败: %s", e, exc_info=True)
            return None

    def _extract_xlsx(self, file_path: Path) -> Optional[str]:
        """提取 Excel 为 Markdown 表格"""
        try:
            from openpyxl import load_workbook
            from openpyxl.utils.exceptions import InvalidFileException
        except ImportError:
            logger.debug("[FileIngestor] openpyxl 未安装")
            return None

        try:
            wb = load_workbook(str(file_path), read_only=True)
            sheets = []
            for ws in wb.worksheets:
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    header = rows[0]
                    separator = "|" + "|".join(["---"] * (len(rows[0].split("|")) - 2)) + "|"
                    sheets.append(f"## {ws.title}\n\n{header}\n{separator}\n" + "\n".join(rows[1:]))
            return "\n\n".join(sheets)
        except (OSError, InvalidFileException, ValueError, TypeError, KeyError) as e:
            logger.warning("[FileIngestor] XLSX 提取失败: %s", e, exc_info=True)
            return None

    def _extract_html(self, file_path: Path) -> Optional[str]:
        """提取 HTML 文本"""
        try:
            from bs4 import BeautifulSoup
            from bs4 import ParserRejectedMarkup
        except ImportError:
            logger.debug("[FileIngestor] beautifulsoup4 未安装")
            return None

        try:
            content = file_path.read_text(encoding="utf-8", errors="replace")
            soup = BeautifulSoup(content, "html.parser")
            # 移除 script 和 style
            for tag in soup(["script", "style"]):
                tag.decompose()
            return soup.get_text(separator="\n", strip=True)
        except (OSError, ParserRejectedMarkup, ValueError, TypeError) as e:
            logger.warning("[FileIngestor] HTML 提取失败: %s", e, exc_info=True)
            return None

    def _extract_epub(self, file_path: Path) -> Optional[str]:
        """提取 epub 文本"""
        try:
            import ebooklib
            from ebooklib import epub
            from ebooklib.epub import EpubException
            from bs4 import BeautifulSoup
            from bs4 import ParserRejectedMarkup
        except ImportError:
            logger.debug("[FileIngestor] ebooklib 未安装")
            return None

        try:
            book = epub.read_epub(str(file_path))
            chapters = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text = soup.get_text(separator="\n", strip=True)
                if text:
                    chapters.append(text)
            return "\n\n".join(chapters)
        except (
            OSError,
            EpubException,
            ParserRejectedMarkup,
            ValueError,
            TypeError,
            KeyError,
        ) as e:
            logger.warning("[FileIngestor] EPUB 提取失败: %s", e, exc_info=True)
            return None

    def _build_file_markdown(self, file_path: Path, text: str) -> str:
        """构建稳定的 canonical Markdown；捕获时间由 raw receipt 单独记录。"""
        return (
            f"# File: {file_path.name}\n\n"
            f"**路径**: `{file_path}`\n"
            f"**类型**: {file_path.suffix.lstrip('.')}\n"
            f"**大小**: {file_path.stat().st_size} bytes\n\n"
            f"---\n\n{text}"
        )

    def _is_supported(self, file_path: Path) -> bool:
        """检查文件类型是否支持"""
        ext = file_path.suffix.lower()
        return (
            ext in _TEXT_EXTENSIONS
            or ext in _PDF_EXTENSIONS
            or ext in _DOCX_EXTENSIONS
            or ext in _PPTX_EXTENSIONS
            or ext in _XLSX_EXTENSIONS
            or ext in _HTML_EXTENSIONS
            or ext in _EPUB_EXTENSIONS
        )

    @staticmethod
    def _now_date() -> str:
        from datetime import datetime

        return datetime.now().strftime("%Y%m%d")
